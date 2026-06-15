from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(0x1E, 0x3A, 0x5F)  # Dark blue
SECONDARY_COLOR = RGBColor(0x2E, 0x86, 0xAB)  # Light blue
ACCENT_COLOR = RGBColor(0xE8, 0x5D, 0x04)  # Orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

# ============================================================
# Slide 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide, PRIMARY_COLOR)

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "🇲🇾 ペナンデジタルノマド", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
    "2ヶ月間 完全計画書", 36, RGBColor(0xFF, 0xD7, 0x00), True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
    "2026年6月16日（火）〜 8月26日（水）", 24, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
    "人数：大人1名 ｜ スタイル：デジタルノマド（単身滞在）", 20, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "📋 目次", 36, PRIMARY_COLOR, True)

toc_items = [
    ("1", "旅費・滞在費 見積もり", "¥400,000（ミドルプラン）"),
    ("2", "2ヶ月間 滞在計画", "週別スケジュール"),
    ("3", "往路ルート", "妙典駅→成田空港→ペナン"),
    ("4", "復路ルート", "ペナン→成田空港→妙典駅"),
    ("5", "成田空港ガイド", "第1ターミナル南ウイング"),
    ("6", "シンガポール乗り継ぎ", "チャンギ空港5h20m待機"),
    ("7", "ペナン到着後", "入国・Grab・宿泊先"),
    ("8", "宿泊先詳細", "Bayan Lepas Airbnb"),
    ("9", "コワーキング・カフェ", "作業環境"),
    ("10", "おすすめ飲食", "ローカルフード"),
    ("11", "観光スポット", "ジョージタウン等"),
    ("12", "持ち物チェックリスト", "パッキングリスト"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.3 + i * 0.48)
    # Number circle
    shape = add_shape(slide, Inches(1), y, Inches(0.4), Inches(0.4), SECONDARY_COLOR)
    shape.text_frame.paragraphs[0].text = num
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, Inches(1.6), y, Inches(4), Inches(0.4),
        title, 16, DARK_GRAY, True)
    # Description
    add_text_box(slide, Inches(6), y, Inches(6), Inches(0.4),
        desc, 14, RGBColor(0x66, 0x66, 0x66))

# ============================================================
# Slide 3: Cost Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "💰 旅費・滞在費 見積もり", 36, PRIMARY_COLOR, True)

# Budget comparison table
budget_data = [
    ["プラン", "総額", "1日あたり", "宿泊レベル"],
    ["🟢 予算（Budget）", "¥250,000", "¥3,470", "シェアハウス"],
    ["🟡 ミドル（Comfort）", "¥400,000", "¥5,560", "Airbnb 1BR"],
    ["🔴 贅沢（Premium）", "¥550,000", "¥7,640", "高級コンド"],
]

# Table
table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.2), Inches(12), Inches(1.5)).table
for i, row_data in enumerate(budget_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = (i == 0)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
                paragraph.font.color.rgb = WHITE
            elif i == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)

# Cost breakdown
cost_items = [
    ("航空券（往復）", "¥63,735"),
    ("宿泊（2ヶ月）", "¥134,696"),
    ("食費（72日）", "¥180,000"),
    ("交通費", "¥14,400"),
    ("SIM/eSIM", "¥1,500"),
    ("海外旅行保険", "¥8,000"),
    ("雑費", "¥21,600"),
    ("予備費（10%）", "¥42,400"),
]

add_text_box(slide, Inches(0.5), Inches(3), Inches(6), Inches(0.5),
    "ミドルプラン 費目別詳細", 20, PRIMARY_COLOR, True)

table = slide.shapes.add_table(len(cost_items) + 1, 2, Inches(0.5), Inches(3.5), Inches(6), Inches(3.5)).table
# Header
table.cell(0, 0).text = "費目"
table.cell(0, 1).text = "金額"
for cell in table.rows[0].cells:
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

for i, (item, cost) in enumerate(cost_items):
    table.cell(i + 1, 0).text = item
    table.cell(i + 1, 1).text = cost

# Total
table.cell(len(cost_items), 0).text = "合計"
table.cell(len(cost_items), 1).text = "¥466,331"
for cell in [table.cell(len(cost_items), 0), table.cell(len(cost_items), 1)]:
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_COLOR
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

# ============================================================
# Slide 4: Schedule Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "📅 2ヶ月間 滞在計画", 36, PRIMARY_COLOR, True)

# Timeline
months = ["6月", "7月", "8月"]
weeks_per_month = [2.5, 4.5, 4]

# Create timeline boxes
timeline_y = Inches(1.5)
box_width = Inches(2.5)

for i, month in enumerate(months):
    x = Inches(0.5 + i * 4)
    add_shape(slide, x, timeline_y, box_width, Inches(0.5), SECONDARY_COLOR)
    add_text_box(slide, x, timeline_y, box_width, Inches(0.5),
        month, 18, WHITE, True, PP_ALIGN.CENTER)

# Week details
weeks = [
    ("第1週", "6/16〜6/22", "到着・セットアップ"),
    ("第2週", "6/23〜6/29", "生活リズム確立"),
    ("第3週", "6/30〜7/6", "ローカル体験"),
    ("第4週", "7/7〜7/13", "観光充実"),
    ("第5週", "7/14〜7/20", "集中作業"),
    ("第6週", "7/21〜7/27", "週末アクティビティ"),
    ("第7週", "7/28〜8/3", "穴場探索"),
    ("第8週", "8/4〜8/10", "お土産探し"),
    ("第9週", "8/11〜8/17", "ラストスパート"),
    ("第10週", "8/18〜8/26", "帰国準備"),
]

add_text_box(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5),
    "週別スケジュール", 20, PRIMARY_COLOR, True)

table = slide.shapes.add_table(11, 3, Inches(0.5), Inches(3), Inches(12), Inches(4)).table
table.cell(0, 0).text = "期間"
table.cell(0, 1).text = "日付"
table.cell(0, 2).text = "内容"
for cell in table.rows[0].cells:
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

for i, (week, dates, content) in enumerate(weeks):
    table.cell(i + 1, 0).text = week
    table.cell(i + 1, 1).text = dates
    table.cell(i + 1, 2).text = content

# ============================================================
# Slide 5: Outbound Route (妙典駅→成田空港)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🚃 往路ルート：妙典駅→成田空港", 36, PRIMARY_COLOR, True)

# Route diagram
route_steps = [
    ("妙典駅（T21）", "05:14発", "東西線 西船橋方面"),
    ("本八幡駅（BZ）", "05:17着", "南口出口"),
    ("京成八幡駅（H07）", "徒歩5分", "約400m"),
    ("成田空港駅（NRT）", "06:42着", "京成本線 約1時間17分"),
]

y = Inches(1.5)
for i, (station, time, note) in enumerate(route_steps):
    x = Inches(1 + i * 3)
    # Station box
    add_shape(slide, x, y, Inches(2.5), Inches(1), SECONDARY_COLOR)
    add_text_box(slide, x, y, Inches(2.5), Inches(0.5),
        station, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.3), Inches(2.5), Inches(0.4),
        time, 12, WHITE, False, PP_ALIGN.CENTER)
    # Arrow
    if i < 3:
        add_text_box(slide, x + Inches(2.5), y + Inches(0.2), Inches(0.5), Inches(0.5),
            "→", 24, ACCENT_COLOR, True, PP_ALIGN.CENTER)

# Notes
add_text_box(slide, Inches(0.5), Inches(3), Inches(12), Inches(0.5),
    "⚠️ 注意事項", 20, ACCENT_COLOR, True)

notes = [
    "• 朝ラッシュ（6:30-8:30）は踏切閉塞1時間中31分間",
    "• 京成本線 片道約¥1,200",
    "• 約1時間28分で成田空港到着",
]

for i, note in enumerate(notes):
    add_text_box(slide, Inches(0.5), Inches(3.5 + i * 0.4), Inches(12), Inches(0.4),
        note, 14, DARK_GRAY)

# ============================================================
# Slide 6: Return Route
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🚃 復路ルート：ペナン→成田空港→妙典駅", 36, PRIMARY_COLOR, True)

# Route diagram
route_steps = [
    ("ペナン空港（PEN）", "17:50発", "スクートTR429"),
    ("シンガポール（SIN）", "19:25着", "乗り継ぎ3時間40分"),
    ("シンガポール（SIN）", "23:05発", "スクートTR884"),
    ("成田空港（NRT）", "07:05+1着", "翌朝到着"),
    ("妙典駅（T21）", "09:00頃", "JR東西線"),
]

y = Inches(1.5)
for i, (station, time, note) in enumerate(route_steps):
    x = Inches(0.5 + i * 2.5)
    # Station box
    add_shape(slide, x, y, Inches(2.2), Inches(1), SECONDARY_COLOR)
    add_text_box(slide, x, y, Inches(2.2), Inches(0.5),
        station, 12, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.3), Inches(2.2), Inches(0.4),
        time, 11, WHITE, False, PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        add_text_box(slide, x + Inches(2.2), y + Inches(0.2), Inches(0.3), Inches(0.5),
            "→", 20, ACCENT_COLOR, True, PP_ALIGN.CENTER)

# Notes
add_text_box(slide, Inches(0.5), Inches(3), Inches(12), Inches(0.5),
    "⚠️ 復路のポイント", 20, ACCENT_COLOR, True)

notes = [
    "• ペナン→シンガポール 1時間35分",
    "• シンガポール空港で3時間40分待機（Jewel散策可能）",
    "• シンガポール→成田 7時間",
    "• 翌朝07:05成田着→JRで妙典駅まで約2時間",
]

for i, note in enumerate(notes):
    add_text_box(slide, Inches(0.5), Inches(3.5 + i * 0.4), Inches(12), Inches(0.4),
        note, 14, DARK_GRAY)

# ============================================================
# Slide 7: Narita Airport Guide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "✈️ 成田空港 第1ターミナル 南ウイング", 36, PRIMARY_COLOR, True)

# Floor map
floors = [
    ("5F", "レストラン・ショップ・展望デッキ"),
    ("4F", "出発ロビー（チェックイン）← スクートは南ウイングIゾーン"),
    ("3F", "出国審査 → 搭乗ゲート"),
    ("1F", "到着ロビー"),
    ("B1F", "JR・京成 成田空港駅"),
]

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
    "フロア構造", 20, PRIMARY_COLOR, True)

table = slide.shapes.add_table(5, 2, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5)).table
for i, (floor, desc) in enumerate(floors):
    table.cell(i, 0).text = floor
    table.cell(i, 1).text = desc
    if i == 1:  # 4F highlighted
        table.cell(i, 0).fill.solid()
        table.cell(i, 0).fill.fore_color.rgb = ACCENT_COLOR
        table.cell(i, 1).fill.solid()
        table.cell(i, 1).fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)

# Check-in flow
add_text_box(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
    "スクート チェックインフロー", 20, PRIMARY_COLOR, True)

flow_steps = [
    "① 4F 南ウイング Iゾーン",
    "② チェックインカウンター",
    "③ 手荷物預け",
    "④ 4F 保安検査場",
    "⑤ 3F 出国審査",
    "⑥ 3F 搭乗ゲート（約15分）",
]

for i, step in enumerate(flow_steps):
    add_text_box(slide, Inches(7), Inches(1.8 + i * 0.4), Inches(6), Inches(0.4),
        step, 14, DARK_GRAY)

# Morning info
add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
    "朝7時台の出国審査", 20, PRIMARY_COLOR, True)

morning_info = [
    "• 手荷物検査：約10分",
    "• 出国審査：約2分",
    "• 搭乗ゲートまで：約15分",
    "• 合計：約30分で搭乗ゲート到着可能",
]

for i, info in enumerate(morning_info):
    add_text_box(slide, Inches(0.5), Inches(5 + i * 0.35), Inches(12), Inches(0.35),
        info, 14, DARK_GRAY)

# ============================================================
# Slide 8: Singapore Transit
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🇸🇬 シンガポール乗り継ぎガイド", 36, PRIMARY_COLOR, True)

# Transit info
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
    "往路（6/16）", 20, PRIMARY_COLOR, True)

transit_info = [
    "• 到着：14:35（チャンギ空港）",
    "• 出発：19:55（同ターミナル）",
    "• 待機時間：5時間20分",
    "• ターミナル：T1",
]

for i, info in enumerate(transit_info):
    add_text_box(slide, Inches(0.5), Inches(1.8 + i * 0.35), Inches(6), Inches(0.35),
        info, 14, DARK_GRAY)

add_text_box(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
    "復路（8/26）", 20, PRIMARY_COLOR, True)

transit_info = [
    "• 到着：19:25（チャンギ空港）",
    "• 出発：23:05（同ターミナル）",
    "• 待機時間：3時間40分",
    "• ターミナル：T1",
]

for i, info in enumerate(transit_info):
    add_text_box(slide, Inches(7), Inches(1.8 + i * 0.35), Inches(6), Inches(0.35),
        info, 14, DARK_GRAY)

# Jewel Changi
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
    "💎 散策スポット：Jewel Changi", 20, PRIMARY_COLOR, True)

jewel_info = [
    "• 世界最大級の室内庭園",
    "• フォール（滝）：高さ40m",
    "• 免税店・レストラン多数",
    "• 入場無料（一部有料施設あり）",
]

for i, info in enumerate(jewel_info):
    add_text_box(slide, Inches(0.5), Inches(4 + i * 0.35), Inches(12), Inches(0.35),
        info, 14, DARK_GRAY)

# Tips
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
    "💡 乗り継ぎのコツ", 20, ACCENT_COLOR, True)

tips = [
    "• SG Arrival Card 事前提出推奨（入国審査スムーズに）",
    "• ターミナル内の電動カート利用可能",
    "• 搭乗開始は出発40分前",
]

for i, tip in enumerate(tips):
    add_text_box(slide, Inches(0.5), Inches(6 + i * 0.35), Inches(12), Inches(0.35),
        tip, 14, DARK_GRAY)

# ============================================================
# Slide 9: Penang Arrival
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🇲🇾 ペナン到着後の流れ", 36, PRIMARY_COLOR, True)

# Arrival steps
steps = [
    ("① 入国審査", "MDAC確認（事前提出済み）\nパスポート提示\n入国スタンプ"),
    ("② 手荷物受取", "カウンター番号確認\n荷物到着待ち"),
    ("③ カウンター", "現地SIM購入可能\n両替コーナー"),
    ("④ Grab配車", "アプリで配車\nRM15-20（¥470-630）"),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.5)
    # Step box
    add_shape(slide, x, y, Inches(2.8), Inches(1.2), SECONDARY_COLOR)
    add_text_box(slide, x, y, Inches(2.8), Inches(0.4),
        title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.4), Inches(2.8), Inches(0.8),
        desc, 11, WHITE, False, PP_ALIGN.CENTER)

# Grab details
add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.5),
    "🚗 Grab バヤン・レパスへのアクセス", 20, PRIMARY_COLOR, True)

grab_info = [
    "• ペナン空港到着ロビー出口で配車",
    "• 料金：RM15〜20（¥470〜630）",
    "• 時間：約15分",
    "• 宿泊先到着後、チェックイン",
]

for i, info in enumerate(grab_info):
    add_text_box(slide, Inches(0.5), Inches(3.7 + i * 0.35), Inches(12), Inches(0.35),
        info, 14, DARK_GRAY)

# MDAC
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
    "📋 MDAC（Malaysia Digital Arrival Card）", 20, ACCENT_COLOR, True)

mdac_info = [
    "• 申請URL：https://imigresen-online.imi.gov.my/mdac/main",
    "• 到着72時間前までに提出必須",
    "• 無料",
    "• パスポート情報・滞在先・渡航目的を入力",
]

for i, info in enumerate(mdac_info):
    add_text_box(slide, Inches(0.5), Inches(6 + i * 0.35), Inches(12), Inches(0.35),
        info, 14, DARK_GRAY)

# ============================================================
# Slide 10: Accommodation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🏠 宿泊先詳細：Bayan Lepas Airbnb", 36, PRIMARY_COLOR, True)

# Property info
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
    "物件情報", 20, PRIMARY_COLOR, True)

property_info = [
    "• 名前：SPICE/FTZマスタールーム、専用トイレ＆駐車場付き",
    "• 場所：マレーシア、ペナン、バヤン・レパス",
    "• 評価：★4.55（レビュー11件）",
    "• チェックイン：2026/6/16",
    "• チェックアウト：2026/8/26",
    "• 人数：1人",
]

for i, info in enumerate(property_info):
    add_text_box(slide, Inches(0.5), Inches(1.8 + i * 0.35), Inches(6), Inches(0.35),
        info, 14, DARK_GRAY)

# Price
add_text_box(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
    "料金", 20, PRIMARY_COLOR, True)

price_info = [
    "• 通常料金：¥128,366/月",
    "• 割引後：¥67,348/月",
    "• 割引率：48% OFF",
    "• 2ヶ月合計：¥134,696",
]

for i, info in enumerate(price_info):
    add_text_box(slide, Inches(7), Inches(1.8 + i * 0.35), Inches(6), Inches(0.35),
        info, 14, DARK_GRAY)

# Features
add_text_box(slide, Inches(0.5), Inches(4), Inches(12), Inches(0.5),
    "設備", 20, PRIMARY_COLOR, True)

features = [
    "✅ 仕事用スペース", "✅ セルフチェックイン", "✅ アイロン",
    "✅ 無料駐車スペース", "✅ プール", "✅ ベッド1個",
    "✅ 個室内のゲスト専用バスルーム", "✅ ペナン国立公園の中",
]

for i, feature in enumerate(features):
    col = i % 4
    row = i // 4
    add_text_box(slide, Inches(0.5 + col * 3), Inches(4.5 + row * 0.4), Inches(3), Inches(0.4),
        feature, 14, DARK_GRAY)

# Host
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
    "ホスト情報", 20, PRIMARY_COLOR, True)

host_info = [
    "• ホスト名：Ckさん",
    "• ホスティング歴：7年",
]

for i, info in enumerate(host_info):
    add_text_box(slide, Inches(0.5), Inches(6 + i * 0.35), Inches(12), Inches(0.35),
        info, 14, DARK_GRAY)

# ============================================================
# Slide 11: Coworking & Cafes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "💻 コワーキング・カフェ", 36, PRIMARY_COLOR, True)

# Coworking spaces
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
    "コワーキングスペース", 20, PRIMARY_COLOR, True)

coworking_data = [
    ["施設名", "Wi-Fi", "料金", "エリア"],
    ["MSOGO", "500Mbps+", "¥3,000/月", "ジョージタウン"],
    ["Common Ground", "300Mbps+", "¥2,500/月", "ジョージタウン"],
    ["The Coffee Commune", "200Mbps+", "ドリンクのみ", "バヤン・レパス"],
]

table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.8), Inches(6), Inches(1.5)).table
for i, row_data in enumerate(coworking_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.bold = True

# Cafes
add_text_box(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
    "おすすめカフェ（作業向き）", 20, PRIMARY_COLOR, True)

cafe_data = [
    ["カフェ名", "Wi-Fi", "電源", "住所"],
    ["China House", "速い", "あり", "ジョージタウン"],
    ["Satchworks", "速い", "あり", "ジョージタウン"],
    ["Macallum Connoisseurs", "速い", "あり", "ジョージタウン"],
]

table = slide.shapes.add_table(4, 4, Inches(7), Inches(1.8), Inches(6), Inches(1.5)).table
for i, row_data in enumerate(cafe_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.bold = True

# Daily routine
add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.5),
    "月〜金のルーティン", 20, PRIMARY_COLOR, True)

routine = [
    "07:00 起床・朝食",
    "08:00-12:00 コワーキング/カフェで作業",
    "12:00-13:00 昼食",
    "13:00-17:00 作業 or 観光",
    "17:00-18:00 urchases（スーパー）",
    "18:00-19:00 夕食",
    "19:00-22:00 自由時間",
    "22:00 就寝",
]

for i, item in enumerate(routine):
    col = i % 4
    row = i // 4
    add_text_box(slide, Inches(0.5 + col * 3), Inches(4.3 + row * 0.4), Inches(3), Inches(0.4),
        item, 14, DARK_GRAY)

# ============================================================
# Slide 12: Food Guide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🍜 おすすめ飲食", 36, PRIMARY_COLOR, True)

# Local food
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
    "ローカルフード（1食RM8〜15 / ¥250〜470）", 20, PRIMARY_COLOR, True)

food_data = [
    ["料理", "値段", "おすすめ店"],
    ["ペナン・ローク mee", "RM8〜12", "平安楼"],
    ["アサム・ラクサ", "RM8〜12", "新峰肉骨茶"],
    ["ナシ・レマ", "RM8〜12", "バーミーコーヒー"],
    ["チキン・リース", "RM10〜15", "ナナバーミーコーヒー"],
    ["ペナン・カレー・ミー", "RM8〜12", "月光バーミーコーヒー"],
]

table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5)).table
for i, row_data in enumerate(food_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.bold = True

# Restaurants
add_text_box(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.5),
    "レストラン（RM20〜50 / ¥630〜1,580）", 20, PRIMARY_COLOR, True)

restaurant_data = [
    ["店名", "料理ジャンル", "特徴"],
    ["竹の子ベトナム料理", "ベトナム", "安くて美味しい"],
    ["太興咖啡室", "カンボジアン", "カフェ併設"],
    ["Red Garden Food Paradise", "フードコート", "多国籍"],
]

table = slide.shapes.add_table(4, 3, Inches(7), Inches(1.8), Inches(6), Inches(1.5)).table
for i, row_data in enumerate(restaurant_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.bold = True

# Tips
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
    "💡 食事のコツ", 20, ACCENT_COLOR, True)

food_tips = [
    "• 朝食はバーミーコーヒーがおすすめ（ローカル定番）",
    "• 夕食はナイトマーケットで Various",
    "• スーパーで食材を購入し、自炊も可能",
    "• タイ料理も人気（タイ国境近く）",
]

for i, tip in enumerate(food_tips):
    add_text_box(slide, Inches(0.5), Inches(5.3 + i * 0.35), Inches(12), Inches(0.35),
        tip, 14, DARK_GRAY)

# ============================================================
# Slide 13: Tourist Spots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🏛️ おすすめ観光スポット", 36, PRIMARY_COLOR, True)

# Spots
spots = [
    ("ジョージタウン", "世界遺産の歴史地区\n美しい壁画が点在\nカフェ・ショップ多数"),
    ("ペナンヒル", "標高833mの山\n標高833mの山\nロープウェイで登山"),
    ("バトゥーフェリンビーチ", "ペナン最大のビーチ\nアクティビティ充実\nサンセットが美しい"),
    ("クライマス洞窟", "仏教寺院\n洞窟の中にある寺院\nスリル満点"),
]

for i, (name, desc) in enumerate(spots):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.5)
    # Spot box
    add_shape(slide, x, y, Inches(2.8), Inches(1.5), SECONDARY_COLOR)
    add_text_box(slide, x, y, Inches(2.8), Inches(0.4),
        name, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.4), Inches(2.8), Inches(1),
        desc, 12, WHITE, False, PP_ALIGN.CENTER)

# Weekly activities
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
    "週末のおすすめアクティビティ", 20, PRIMARY_COLOR, True)

weekly_activities = [
    ("第1週", "ジョージタウン壁画巡り、ペナンヒル"),
    ("第2週", "バトゥーフェリンビーチ、海上マスジット"),
    ("第3週", "クライマス洞窟、ゴッテンテンプル"),
    ("第4週", "ペナン国立公園トレイル、バターワース散策"),
    ("第5週", "ケッピングNature Reserve、マーライ・パダー"),
    ("第6週", "ペナン・バトル・トラック、アーケードゲーム"),
    ("第7週", "タイ佛教寺院、中国寺院巡り"),
    ("第8週", "ペナン浮き市場、ローカルフードツアー"),
    ("第9週", "お土産探し、ラスト観光"),
    ("第10週", "帰国準備"),
]

table = slide.shapes.add_table(11, 2, Inches(0.5), Inches(4), Inches(12), Inches(3)).table
table.cell(0, 0).text = "期間"
table.cell(0, 1).text = "おすすめアクティビティ"
for cell in table.rows[0].cells:
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

for i, (week, activity) in enumerate(weekly_activities):
    table.cell(i + 1, 0).text = week
    table.cell(i + 1, 1).text = activity

# ============================================================
# Slide 14: Packing List
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "🧳 持ち物チェックリスト", 36, PRIMARY_COLOR, True)

# Packing categories
categories = [
    ("衣服", [
        "Tシャツ（7枚）",
        "半ズボン（4本）",
        "長ズボン（1本）",
        "羽織り（冷房対策）",
        "水着",
        "サンダル",
        "スニーカー",
    ]),
    ("貴重品", [
        "パスポート（コピーも携帯）",
        "航空券控え",
        "保険証券",
        "クレジットカード（2枚以上）",
        "現金（日本円・マレーシアリンギット）",
    ]),
    ("電子機器", [
        "ノートPC＋充電器",
        "スマホ＋充電ケーブル",
        "モバイルバッテリー（10,000mAh〜）",
        "変換プラグ（Gタイプ）",
    ]),
    ("日用品", [
        "常備薬",
        "日焼け止め（SPF50+）",
        "虫除けスプレー",
        "折りたたみ傘",
    ]),
]

for i, (category, items) in enumerate(categories):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.5)
    # Category header
    add_shape(slide, x, y, Inches(2.8), Inches(0.5), PRIMARY_COLOR)
    add_text_box(slide, x, y, Inches(2.8), Inches(0.5),
        category, 14, WHITE, True, PP_ALIGN.CENTER)
    # Items
    for j, item in enumerate(items):
        add_text_box(slide, x, y + Inches(0.5 + j * 0.35), Inches(2.8), Inches(0.35),
            f"□ {item}", 12, DARK_GRAY)

# Checklist
add_text_box(slide, Inches(0.5), Inches(5), Inches(12), Inches(0.5),
    "📋 出発前チェックリスト", 20, ACCENT_COLOR, True)

checklist = [
    "□ パスポート（残存期間確認済み）",
    "□ 航空券（控え／スマホ画面可）",
    "□ 海外旅行保険証券",
    "□ クレジットカード（2枚以上推奨）",
    "□ 現金（日本円・マレーシアリンギット）",
    "□ ノートPC + 充電器",
    "□ スマホ + 充電ケーブル",
    "□ モバイルバッテリー",
    "□ 変換プラグ（マレーシアはGタイプ）",
    "□ 常備薬",
    "□ 日焼け止め・虫除け",
    "□ 折りたたみ傘（雨季です）",
]

for i, item in enumerate(checklist):
    col = i % 3
    row = i // 3
    add_text_box(slide, Inches(0.5 + col * 4), Inches(5.5 + row * 0.35), Inches(4), Inches(0.35),
        item, 12, DARK_GRAY)

# ============================================================
# Slide 15: Final Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, PRIMARY_COLOR)

add_text_box(slide, Inches(1), Inches(1), Inches(11), Inches(1),
    "🇲🇾 ペナンデジタルノマド 2ヶ月間", 36, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
    "2026年6月16日（火）〜 8月26日（水）", 24, RGBColor(0xFF, 0xD7, 0x00), True, PP_ALIGN.CENTER)

summary_items = [
    "✈️ 往復航空券：¥63,735（スクート）",
    "🏠 宿泊：¥134,696（2ヶ月）",
    "💰 総額：約¥400,000（ミドルプラン）",
    "📍 場所：バヤン・レパス、ペナン",
]

for i, item in enumerate(summary_items):
    add_text_box(slide, Inches(1), Inches(3.5 + i * 0.5), Inches(11), Inches(0.5),
        item, 20, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6), Inches(11), Inches(0.8),
    "楽しいペナン生活を！", 28, RGBColor(0xFF, 0xD7, 0x00), True, PP_ALIGN.CENTER)

# Save presentation
output_path = r"C:\Users\GoldRush\Documents\MyProject\AsiaJurney\ペナン2ヶ月間_旅のしおり.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
